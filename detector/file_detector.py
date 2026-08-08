import os
import fitz
from docx import Document
from pptx import Presentation


def extract_text_from_file(file_path, password=None):

    extension = os.path.splitext(
        file_path
    )[1].lower()

    text = ""

    # =====================
    # TXT
    # =====================

    if extension == ".txt":

        with open(

            file_path,

            "r",

            encoding="utf-8",

            errors="ignore"

        ) as file:

            text = file.read()

    # =====================
    # PDF (PyMuPDF)
    # =====================

    elif extension == ".pdf":

        pdf = fitz.open(file_path)

        if pdf.is_encrypted:

            if password is None:

                pdf.close()

                raise Exception(
                    "Password required."
                )

            authenticated = pdf.authenticate(password)

            if not authenticated:

                pdf.close()

                raise Exception(
                    "Incorrect password."
                )

        for page in pdf:

            page_text = page.get_text()

            if page_text:

                text += page_text + "\n"

        pdf.close()

    # =====================
    # DOCX
    # =====================

    elif extension == ".docx":

        document = Document(file_path)

        for paragraph in document.paragraphs:

            text += paragraph.text + "\n"

    # =====================
    # PPTX
    # =====================

    elif extension == ".pptx":

        presentation = Presentation(file_path)

        for slide in presentation.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    text += shape.text + "\n"

    return text.strip()